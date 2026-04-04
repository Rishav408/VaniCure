from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

prs = Presentation()

def set_title_format(title_shape):
    for paragraph in title_shape.text_frame.paragraphs:
        paragraph.font.name = 'Arial'
        paragraph.font.bold = True

def add_title_slide(title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = slide.shapes.title
    title_shape.text = title
    set_title_format(title_shape)
    
    subtitle_shape = slide.placeholders[1]
    subtitle_shape.text = subtitle

def add_bullet_slide(title, points):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_shape = slide.shapes.title
    title_shape.text = title
    set_title_format(title_shape)
    
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    
    if points:
        tf.text = points[0]
        for point in points[1:]:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0
            p.font.name = 'Arial'
            p.font.size = Pt(20)
            
        tf.paragraphs[0].font.name = 'Arial'
        tf.paragraphs[0].font.size = Pt(20)

def add_image_slide(title, image_path):
    slide = prs.slides.add_slide(prs.slide_layouts[5]) 
    title_shape = slide.shapes.title
    title_shape.text = title
    set_title_format(title_shape)
    
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(1), Inches(1.5), width=Inches(8))

# --- Custom Architecture Slides (Native PPT Shapes so Canva imports them perfectly) ---

def add_architecture_parallel(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "System Architecture: Parallel Execution (Production Defaults ~52ms)"
    set_title_format(title)
    
    shapes = slide.shapes
    
    # 1. React Client
    left = Inches(0.4)
    top = Inches(3.5)
    width = Inches(2.0)
    height = Inches(1.2)
    s1 = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s1.text = "1. React Frontend\n(Live Mic/Audio Blob)"
    
    # Arrow
    shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left + width + Inches(0.1), top + Inches(0.4), Inches(0.4), Inches(0.4))
    
    # 2. API Threadpool
    left += width + Inches(0.6)
    s2 = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s2.fill.solid()
    s2.fill.fore_color.rgb = RGBColor(0, 102, 204)
    s2.text = "2. FastAPI Edge Server\n(ThreadPoolExecutor)\nasyncio.gather()"
    
    # 3. Pipeline A (Top)
    left_pipe = left + width + Inches(0.8)
    s3 = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pipe, top - Inches(1.5), Inches(2.2), Inches(1.4))
    s3.fill.solid()
    s3.fill.fore_color.rgb = RGBColor(102, 51, 153)
    s3.text = "Pipeline A (Parallel)\n• PANNs CNN14\n• YAMNet\n• CNN-BiLSTM"
    
    # 3. Pipeline B (Bottom)
    s4 = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pipe, top + Inches(1.3), Inches(2.2), Inches(1.4))
    s4.fill.solid()
    s4.fill.fore_color.rgb = RGBColor(102, 51, 153)
    s4.text = "Pipeline B (Parallel)\n• ResNet18\n• MobileNetV2\n• TinyCNN"
    
    # 4. Pipeline C (Combine)
    left_end = left_pipe + Inches(2.2) + Inches(0.6)
    s5 = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_end, top, Inches(1.8), height)
    s5.fill.solid()
    s5.fill.fore_color.rgb = RGBColor(204, 0, 0)
    s5.text = "Pipeline C\n(Final Risk Average)"

def add_architecture_sequential(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "System Architecture: Sequential Execution (Hardware Throttling ~76ms)"
    set_title_format(title)
    
    shapes = slide.shapes
    
    # 1. API
    left = Inches(0.5)
    top = Inches(3.5)
    width = Inches(2.0)
    height = Inches(1.4)
    s1 = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s1.fill.solid()
    s1.fill.fore_color.rgb = RGBColor(0, 102, 204)
    s1.text = "1. FastAPI Server\n(Sequential Mode Triggered)"
    
    # Arrow
    shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left + width + Inches(0.1), top + Inches(0.5), Inches(0.4), Inches(0.4))
    
    # 2. Block 1 (Pipeline A)
    left += width + Inches(0.6)
    s2 = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s2.fill.solid()
    s2.fill.fore_color.rgb = RGBColor(102, 51, 153)
    s2.text = "2. Wait Block 1\nExecute Pipeline A\n(PANNs, YAMNet, BiLSTM)"
    
    # Arrow
    shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left + width + Inches(0.1), top + Inches(0.5), Inches(0.4), Inches(0.4))

    # 3. Block 2 (Pipeline B)
    left += width + Inches(0.6)
    s3 = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s3.fill.solid()
    s3.fill.fore_color.rgb = RGBColor(102, 51, 153)
    s3.text = "3. Wait Block 2\nExecute Pipeline B\n(ResNet, MobileNet, TinyCNN)"

    # Arrow
    shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left + width + Inches(0.1), top + Inches(0.5), Inches(0.4), Inches(0.4))

    # 4. Pipeline C
    left += width + Inches(0.6)
    width = Inches(1.8)
    s4 = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s4.fill.solid()
    s4.fill.fore_color.rgb = RGBColor(204, 0, 0)
    s4.text = "4. Final Output\nAggregate Pipeline C"

# --- Presentation Content Gen ---

add_title_slide(
    "VaniCure: AI-Powered Respiratory Diagnostics", 
    "Offline AI Diagnostic Agent for TB & Asthma Screening\n\nExecution Strategy: Edge-Device Parallel & Sequential Inference"
)

add_bullet_slide("1. Introduction & Motivation", [
    "Problem: Rural Primary Health Centres (PHCs) lack X-rays and reliable clinical diagnostics, with very low internet connectivity.",
    "Goal: Build an offline-first AI system that analyzes patient respiratory audio (coughs, breathing) to predict Tuberculosis and Asthma risk.",
    "Advantage: Complete privacy via on-device inference (Zero Cloud dependency) and non-invasive microphone-based screening.",
    "Engineering Challenge: Running multiple Deep Learning networks simultaneously on standard hardware severely bottlenecks CPUs."
])

add_bullet_slide("2. Project Objectives", [
    "Develop a robust 6-model AI ensemble utilizing diverse structures (CNNs, LSTMs, MobileNets).",
    "Organize models hierarchically into 3 distinct diagnostic pipelines to prevent individual model hallucination.",
    "Engineer an Asynchronous Execution Engine managing both Sequential and Parallel threading.",
    "Establish a full-stack React frontend & FastAPI edge application.",
    "Achieve total inference times under < 100ms via optimized concurrency."
])

# INJECT NATIVE ARCHITECTURE SLIDES
add_architecture_parallel(prs)
add_architecture_sequential(prs)

add_bullet_slide("5. AI Model Architectures (6 Models)", [
    "Pipeline A (Heavy Pretrained Extractors):",
    "  • PANNs CNN14: 14-layer CNN mapping massive 64-bin Mel-Spectrogram features.",
    "  • YAMNet: Google’s lightweight MobileNet v1 producing rapid 521-class logits.",
    "  • CNN-BiLSTM: Custom-trained encoding sequence for cough space + breathing rhythm.",
    "",
    "Pipeline B (Lightweight Edge-Optimized Cast):",
    "  • ResNet18-Audio: Processes residual skip-connections across spectral boundaries.",
    "  • MobileNetV2-Audio: Shrinks computing layers using depthwise separable convolutions.",
    "  • TinyCNN: 3-Layer ultra-fast baseline network executing under <20ms."
])

add_bullet_slide("6. Performance Benchmark Metrics", [
    "Pipeline Accuracy Hierarchy:",
    "  • Pipeline A Absolute Mean: 84.2%",
    "  • Pipeline B Absolute Mean: 81.2%",
    "  • Pipeline C Combined Output: 82.7%",
    "",
    "System Sub-Latency Matrix:",
    "  • Strict Sequential Pipeline Pathing: ~76 ms",
    "  • Fully Unlocked Parallel Execution: ~52 ms (2.3x Benchmark Speedup vs singular execution)",
    "  • Result: Seamless instantaneous clinical feedback."
])

add_image_slide(
    "7. Live Architecture: Parallel Execution Visualized",
    "/Users/kunaltailor/.gemini/antigravity/brain/40cb226d-092d-46b3-80e9-95813d5a78a1/pipeline_cards.png"
)

add_bullet_slide("8. Risk Assessment Algorithm", [
    "The Pipeline C combined aggregate dictates final action parameters:",
    "",
    "🔴 CRITICAL (Suspect TB):",
    "  • Trigger: Pipeline C TB Risk > 25%",
    "  • Direction: Immediate referral for sputum/GeneXpert examination.",
    "",
    "🟡 WARNING (COPD / Asthma Variant):",
    "  • Trigger: Pipeline C Asthma Risk > 30%",
    "  • Direction: Pulmonology consult recommended within 48-hours.",
    "",
    "🟢 NORMAL / LOW RISK:",
    "  • Trigger: All levels beneath thresholds.",
    "  • Direction: General 3-month clinic checkup suggested."
])

add_image_slide(
    "9. Local Output: Outbreak District Dashboard",
    "/Users/kunaltailor/.gemini/antigravity/brain/40cb226d-092d-46b3-80e9-95813d5a78a1/dashboard.png"
)

add_bullet_slide("10. Future Scope & Conclusions", [
    "Conclusion:",
    "  • VaniCure proves that advanced 6-model ML workflows can run offline inside low-connectivity clinical edge hardware.",
    "  • Combining Sequential fallback strategies prevents system crashes, while Parallel unlocking enables state-of-the-art speeds.",
    "",
    "Future Scope:",
    "  • Hardware Recognition: Auto-switching execution nodes based on RAM capacity on boot.",
    "  • Clinical Trials: Gathering metrics beside Gold Standard hospital assays.",
    "  • Edge Containers: Compiling models directly to native mobile TF-Lite architecture."
])

output_path = "/Users/kunaltailor/Desktop/Kunal/VaniCure/VaniCure/VaniCure_Presentation.pptx"
prs.save(output_path)
print(f"Presentation successfully saved to: {output_path}")
